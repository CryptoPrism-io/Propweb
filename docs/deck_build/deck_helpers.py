"""Reusable primitives for building PropWeb deck slides that match the
existing teal+sand / Calibri visual grammar (see
docs/superpowers/specs/2026-07-14-pitch-deck-design.md section 1)."""

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

COLORS = {
    'teal': '0E6E5C',
    'deep_teal': '0A4A3E',
    'sand': 'F2F7F5',
    'sand_line': 'D8DEE4',
    'graphite': '141821',
    'grey': '5A6270',
    'amber': 'E8A13D',
    'white': 'FFFFFF',
    'grey_light': 'C9D2D9',
    'grey_dark': '8A94A0',
}

ALIGN = {'l': PP_ALIGN.LEFT, 'ctr': PP_ALIGN.CENTER, 'r': PP_ALIGN.RIGHT}


def hex_color(name_or_hex):
    return RGBColor.from_string(COLORS.get(name_or_hex, name_or_hex))


def add_textbox(slide, x, y, cx, cy, text, size_pt=12.5, bold=False, color='graphite',
                 align='l', anchor='ctr', spc=None, font='Calibri'):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(cx), Emu(cy))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if anchor == 'ctr' else MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = ALIGN[align]
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_color(color)
    if spc is not None:
        run._r.get_or_add_rPr().set('spc', str(spc))
    return tb


def add_mixed_textbox(slide, x, y, cx, cy, runs, size_pt=14, align='l', anchor='ctr', font='Calibri'):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(cx), Emu(cy))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if anchor == 'ctr' else MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = ALIGN[align]
    for r in runs:
        run = p.add_run()
        run.text = r['text']
        run.font.name = font
        run.font.size = Pt(r.get('size_pt', size_pt))
        run.font.bold = r.get('bold', False)
        run.font.color.rgb = hex_color(r.get('color', 'graphite'))
    return tb


def add_card(slide, x, y, cx, cy, fill='sand', border='sand_line', adj_raw=6061, border_w_emu=9525):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(y), Emu(cx), Emu(cy))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_color(fill)
    shape.line.color.rgb = hex_color(border)
    shape.line.width = Emu(border_w_emu)
    shape.shadow.inherit = False
    avLst = shape._element.spPr.find(qn('a:prstGeom')).find(qn('a:avLst'))
    gd = avLst.find(qn('a:gd'))
    if gd is None:
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
    gd.set('fmla', f'val {adj_raw}')
    return shape


def add_stat_row(slide, stats, y=2103120, height=2377440):
    n = len(stats)
    margin = 640080
    gap = 209832
    card_w = (12192000 - 2 * margin - gap * (n - 1)) // n
    for i, s in enumerate(stats):
        x = margin + i * (card_w + gap)
        add_card(slide, x, y, card_w, height, fill='sand', border='sand_line', adj_raw=3077)
        add_textbox(slide, x, y + 274320, card_w, 868680, s['number'],
                    size_pt=36, bold=True, color='teal', align='ctr')
        add_textbox(slide, x + 182880, y + 1188720, card_w - 365760, 1051560, s['label'],
                    size_pt=12.5, color='grey', align='ctr', anchor='top')


def set_dark_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_color('graphite')


def set_light_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_color('white')


def add_blank_slide(prs, light=True):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    (set_light_background if light else set_dark_background)(slide)
    return slide


def add_header(slide, eyebrow_text, h1_text):
    add_textbox(slide, 640080, 502920, 7315200, 320040, eyebrow_text,
                size_pt=12, bold=True, color='teal', spc=300)
    add_textbox(slide, 640080, 868680, 10881360, 777240, h1_text,
                size_pt=30, bold=True, color='graphite')


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def get_slide_text(slide):
    chunks = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    chunks.append(r.text)
        elif shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            chunks.append(r.text)
    return ' '.join(chunks)


def get_notes_text(slide):
    if not slide.has_notes_slide:
        return ''
    return slide.notes_slide.notes_text_frame.text


def add_table(slide, x, y, cx, cy, rows, col_widths):
    n_rows, n_cols = len(rows), len(rows[0])
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(y), Emu(cx), Emu(cy))
    table = graphic_frame.table
    for c, w in enumerate(col_widths):
        table.columns[c].width = Emu(w)
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_color('sand' if (r == 0 or r % 2 == 0) else 'white')
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = 'Calibri'
            run.font.size = Pt(12 if r else 12.5)
            run.font.bold = (r == 0)
            run.font.color.rgb = hex_color('graphite' if r == 0 else 'grey')
    return table


def remove_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slide_id_elements = list(xml_slides)
    r_id = slide_id_elements[index].get(qn('r:id'))
    prs.part.drop_rel(r_id)
    xml_slides.remove(slide_id_elements[index])


def remove_slide_by_text(prs, substr):
    slides_list = list(prs.slides)
    for i, slide in enumerate(slides_list):
        if substr in get_slide_text(slide):
            remove_slide(prs, i)
            return
    raise ValueError(f'no slide found containing {substr!r}')
