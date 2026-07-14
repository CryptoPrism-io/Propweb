"""Verification checks for the pitch deck build. Run standalone:
    python docs/deck_build/verify_deck.py
Each check_taskN function is appended by its corresponding plan task."""

from pptx import Presentation
from deck_helpers import get_slide_text, get_notes_text

SOURCE = 'source.pptx'
OUTPUT = '../PropWeb_Pitch_v2.pptx'


def check(condition, message):
    status = 'PASS' if condition else 'FAIL'
    print(f'[{status}] {message}')
    if not condition:
        raise AssertionError(message)


def find_slide_by_text(prs, substr):
    for i, slide in enumerate(prs.slides):
        if substr in get_slide_text(slide):
            return i, slide
    return None, None


def check_task1():
    prs = Presentation(SOURCE)
    check(len(prs.slides) == 9, 'source.pptx has 9 slides')
    i, _ = find_slide_by_text(prs, 'PROPWEB')
    check(i == 0, 'source.pptx slide 1 contains PROPWEB title')

    out = Presentation(OUTPUT)
    check(len(out.slides) >= 9, 'output has at least the 9 source slides')


EXISTING_NOTES = {
    'PROPWEB': "Open cold: this is the one-line pitch. Say it, don't read it — \"PropWeb is India's trust-first, AI-powered rental platform: verified owners meet verified tenants directly, no brokers, no fake listings, no spam calls.\" Pause after the tagline before moving to the problem slide — let the promise land before you justify it.",
    'Renting in India is broken': "Walk both columns — don't just read the tenant side. The room already knows renting is broken; the point is naming both sides of the pain so the solution slide lands as a genuine two-sided fix, not a tenant-only app. If asked \"isn't this just NoBroker,\" hold that for the next slide.",
    'The hidden cost is trust': "These four numbers carry the slide — let them breathe, don't rush to the NoBroker line. The NoBroker callout is the pre-emptive answer to \"hasn't this already been solved\" — say revenue and loss together, in the same breath: proof of demand, proof the trust layer is still broken.",
    'A trust-first rental platform with an AI engine': "Six pillars, but only two are truly new to the room: Pay-to-Connect and the Trust Score — spend your time there. Verified badges and matchmaking are easy to nod along to; the pay-to-connect economics are what actually kills fake enquiries, so be ready to defend the fee-amount question.",
    'From verified profile to matched move-in': "This is the natural hand-off to the live demo — after \"Connect,\" say \"and I can show you this right now\" and switch to the demo app. Have it already loaded to the Koramangala search before this slide so the transition has no dead air.",
    'Problems no one in India solves well': "This slide is doing double duty as the moat argument, along with the closing slide — scam prevention, background checks and digital agreements are what a trust-first platform can sell that a lead-selling portal structurally can't retrofit. If a CEO asks \"why can't NoBroker just copy this,\" point back here.",
    'Two engines: connect today, services tomorrow': "Two engines, two timeframes — be explicit that \"today\" funds the company and \"tomorrow\" is where the real upside is, using NoBroker's own ~50% non-listing mix as the proof point. Don't let this slide turn into a pricing negotiation; final pricing is explicitly TBD with the CEOs.",
    'Existing portals sell leads': "This is the mic-drop — deliver the two lines slowly, then stop talking. Let \"Existing portals sell leads. PropWeb sells trust.\" sit before you open the floor. If the ask & close slide raised a decision, circle back to it once more before Q&A.",
}


def check_task2():
    prs = Presentation(OUTPUT)
    for marker, expected_note in EXISTING_NOTES.items():
        i, slide = find_slide_by_text(prs, marker)
        check(slide is not None, f'slide containing {marker!r} exists')
        notes = get_notes_text(slide)
        check(len(notes) > 80, f'slide {marker!r} has a real presenter note (got {len(notes)} chars)')
        check(notes == expected_note, f'slide {marker!r} note text matches exactly')


def check_task3():
    prs = Presentation(OUTPUT)
    i, slide = find_slide_by_text(prs, 'A large market with an unsolved trust gap')
    check(slide is not None, 'Market opportunity slide exists')
    text = get_slide_text(slide)
    for expected in ('THE OPPORTUNITY', '1.3–1.7B', '12–19%', '$20B', 'Bengaluru', '803 Cr'):
        check(expected in text, f'Market slide contains {expected!r}')
    check(len(get_notes_text(slide)) > 80, 'Market slide has presenter notes')


def check_task4():
    prs = Presentation(OUTPUT)
    i, slide = find_slide_by_text(prs, 'One well-organised building, not eleven services')
    check(slide is not None, 'Architecture slide exists')
    text = get_slide_text(slide)
    for expected in ('ARCHITECTURE', 'modular monolith', 'Listings & Search', 'Trust & KYC',
                      'Matching Engine', 'Connect & Payments', 'TypeScript'):
        check(expected in text, f'Architecture slide contains {expected!r}')
    check(len(get_notes_text(slide)) > 80, 'Architecture slide has presenter notes')


def check_task5():
    prs = Presentation(OUTPUT)
    matches = [i for i, s in enumerate(prs.slides)
               if 'A modern stack, corrected for India economics' in get_slide_text(s)]
    check(len(matches) == 1, f'exactly one Stack slide exists (found {len(matches)})')
    slide = list(prs.slides)[matches[0]]
    text = get_slide_text(slide)
    for expected in ('THE STACK', 'Next.js', 'Typesense', 'Algolia', 'Leegality', 'DigiLocker', 'AWS Mumbai'):
        check(expected in text, f'Stack slide contains {expected!r}')
    check('12-month build budget' not in text, 'old budget line removed from Stack slide (budget has its own slide now)')
    check(len(get_notes_text(slide)) > 80, 'Stack slide has presenter notes')


CHECKS = [check_task1, check_task2, check_task3, check_task4, check_task5]

if __name__ == '__main__':
    for fn in CHECKS:
        fn()
    print(f'\n{len(CHECKS)} check group(s) passed.')
